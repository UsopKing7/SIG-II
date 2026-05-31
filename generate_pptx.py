#!/usr/bin/env python3
"""Genera presentación PowerPoint: SIG Teleférico Universitario"""

import sys

# ── Add wheel paths ──
sys.path.insert(0, open('/tmp/pptx_path.txt').read().strip())
sys.path.insert(0, '/tmp/py_deps2')  # lxml compiled for cp314
sys.path.insert(0, '/tmp/py_deps3')  # Pillow compiled for cp314

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Colors ──
BG_DARK = RGBColor(0x0F, 0x17, 0x2A)
BLUE = RGBColor(0x3B, 0x82, 0xF6)
GREEN = RGBColor(0x10, 0xB9, 0x81)
WHITE = RGBColor(0xF1, 0xF5, 0xF9)
GRAY = RGBColor(0x94, 0xA3, 0xB8)
DARK_GRAY = RGBColor(0x47, 0x56, 0x69)
CARD_BG = RGBColor(0x1E, 0x29, 0x3B)
CARD_BORDER = RGBColor(0x33, 0x41, 0x55)
LIGHT_TEXT = RGBColor(0xCB, 0xD5, 0xE1)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H


def dark_bg(slide):
    """Set dark background."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BG_DARK


def add_shape(slide, left, top, width, height, fill_color=None, line_color=None, radius=None):
    """Add a rounded rectangle shape."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color or CARD_BG
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    # Adjust corner radius
    shape.adjustments[0] = 0.05
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=LIGHT_TEXT, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name='Calibri'):
    """Add a text box."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_box(slide, left, top, width, height, items, font_size=16,
                   color=LIGHT_TEXT, title=None, title_color=WHITE):
    """Add a text box with bullet points."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    if title:
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(font_size + 2)
        p.font.bold = True
        p.font.color.rgb = title_color
        p.font.name = 'Calibri'
        p.space_after = Pt(8)
        for item in items:
            p = tf.add_paragraph()
            p.text = f"  • {item}"
            p.font.size = Pt(font_size)
            p.font.color.rgb = color
            p.font.name = 'Calibri'
            p.space_after = Pt(4)
    else:
        for i, item in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = f"  • {item}"
            p.font.size = Pt(font_size)
            p.font.color.rgb = color
            p.font.name = 'Calibri'
            p.space_after = Pt(4)
    return txBox


def add_title_bar(slide, icon, title, color=BLUE):
    """Add a styled title bar at the top."""
    # Accent line
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, Pt(4)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()

    # Title text
    txt = f"{icon}  {title}"
    add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.8),
                 txt, font_size=28, color=WHITE, bold=True)


def add_card(slide, left, top, width, height, icon, title, body,
             border_color=BLUE):
    """Add a card-style content box."""
    shape = add_shape(slide, left, top, width, height,
                      fill_color=CARD_BG, line_color=CARD_BORDER)

    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(12)
    tf.margin_right = Pt(12)
    tf.margin_top = Pt(16)
    tf.margin_bottom = Pt(12)

    # Icon + Title
    p = tf.paragraphs[0]
    p.text = f"{icon}  {title}"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = border_color
    p.font.name = 'Calibri'
    p.space_after = Pt(8)

    # Body
    p = tf.add_paragraph()
    p.text = body
    p.font.size = Pt(14)
    p.font.color.rgb = LIGHT_TEXT
    p.font.name = 'Calibri'
    p.space_after = Pt(4)

    return shape


def add_divider(slide, left, top, width, color=BLUE):
    """Add a thin colored divider line."""
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, Pt(3)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()
    return line


# ════════════════════════════════════════════════════════════
# SLIDE 1 — PORTADA
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
dark_bg(slide)

# Gradient-like effect: blue top bar
bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, Inches(3.2))
bar.fill.solid()
bar.fill.fore_color.rgb = RGBColor(0x12, 0x1D, 0x3A)
bar.line.fill.background()

# Diagonal accent shape
accent = slide.shapes.add_shape(
    MSO_SHAPE.RIGHT_TRIANGLE, Inches(9.5), Inches(0), Inches(4), Inches(3.5)
)
accent.fill.solid()
accent.fill.fore_color.rgb = RGBColor(0x14, 0x25, 0x45)
accent.line.fill.background()
accent.rotation = 0.0

# Small accent triangle
accent2 = slide.shapes.add_shape(
    MSO_SHAPE.RIGHT_TRIANGLE, Inches(0), Inches(6.2), Inches(3), Inches(1.8)
)
accent2.fill.solid()
accent2.fill.fore_color.rgb = RGBColor(0x14, 0x25, 0x45)
accent2.line.fill.background()

# University label
add_text_box(slide, Inches(0.8), Inches(0.5), Inches(6), Inches(0.5),
             "UNIVERSIDAD", font_size=12, color=DARK_GRAY, bold=True)

# Main title
add_text_box(slide, Inches(0.8), Inches(1.1), Inches(8), Inches(1.2),
             "SIG Teleférico Universitario", font_size=44, color=WHITE, bold=True)

# Subtitle
add_text_box(slide, Inches(0.8), Inches(2.3), Inches(8), Inches(0.7),
             "Sistema de Información Geográfica para red de transporte universitario",
             font_size=16, color=GRAY)

# Divider
add_divider(slide, Inches(0.8), Inches(3.0), Inches(2), BLUE)

# Author info
add_text_box(slide, Inches(0.8), Inches(3.8), Inches(4), Inches(0.5),
             "Integrantes:", font_size=14, color=DARK_GRAY, bold=True)
add_text_box(slide, Inches(0.8), Inches(4.3), Inches(4), Inches(0.4),
             "  👤  Integrante 1", font_size=16, color=LIGHT_TEXT)
add_text_box(slide, Inches(0.8), Inches(4.8), Inches(4), Inches(0.4),
             "  👤  Integrante 2", font_size=16, color=LIGHT_TEXT)

# Year
add_text_box(slide, Inches(0.8), Inches(5.6), Inches(4), Inches(0.4),
             "📅  2026", font_size=12, color=DARK_GRAY)

# Right decorative icon
add_text_box(slide, Inches(9.5), Inches(4.0), Inches(3), Inches(1.5),
             "🗺️", font_size=64, color=RGBColor(0x1E, 0x29, 0x3B),
             alignment=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════
# SLIDE 2 — PROBLEMÁTICA
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(slide)
add_title_bar(slide, "⚠️", "Problemática", BLUE)

cards_data = [
    ("🗺️", "Visualización georreferenciada",
     "Falta de un mapa unificado que muestre la red de teleférico con datos espaciales.", BLUE),
    ("🔗", "Gestión centralizada",
     "Dificultad para administrar líneas, estaciones, cables y postes desde un solo lugar.", GREEN),
    ("💾", "Sistema CRUD espacial",
     "No existe un sistema con operaciones CRUD que maneje datos con geometrías PostGIS.", BLUE),
]
card_w = Inches(3.5)
card_h = Inches(3.2)
gap = Inches(0.4)
total_w = 3 * card_w + 2 * gap
start_x = (SLIDE_W - total_w) // 2

for i, (icon, title, body, color) in enumerate(cards_data):
    x = start_x + i * (card_w + gap)
    add_card(slide, x, Inches(1.8), card_w, card_h, icon, title, body, color)

add_text_box(slide, Inches(0.8), Inches(5.5), Inches(11), Inches(0.5),
             "💡  La información de la infraestructura está dispersa y no hay un mapa unificado.",
             font_size=14, color=DARK_GRAY, alignment=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════
# SLIDE 3 — OBJETIVOS
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(slide)
add_title_bar(slide, "🎯", "Objetivos", GREEN)

# General objective card
add_card(slide, Inches(1.2), Inches(1.5), Inches(10.5), Inches(1.8),
         "🎯", "Objetivo General",
         "Desarrollar un SIG web para visualizar y gestionar la red de teleférico universitario.",
         GREEN)

# Specific objectives
specs = [
    ("📐", "Modelar geometrías espaciales con PostGIS", BLUE),
    ("☁️", "Exponer API REST con NestJS", GREEN),
    ("🗺️", "Visualizar el mapa 3D con Mapbox GL", BLUE),
]
spec_w = Inches(3.2)
spec_h = Inches(2.2)
spec_gap = Inches(0.35)
spec_total = 3 * spec_w + 2 * spec_gap
spec_start = (SLIDE_W - spec_total) // 2

for i, (icon, title, color) in enumerate(specs):
    x = spec_start + i * (spec_w + spec_gap)
    shape = add_shape(slide, x, Inches(3.8), spec_w, spec_h,
                      fill_color=CARD_BG, line_color=CARD_BORDER)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(12)
    tf.margin_right = Pt(12)
    tf.margin_top = Pt(20)
    p = tf.paragraphs[0]
    p.text = icon
    p.font.size = Pt(28)
    p.font.color.rgb = color
    p.alignment = PP_ALIGN.CENTER
    p = tf.add_paragraph()
    p.text = title
    p.font.size = Pt(13)
    p.font.color.rgb = LIGHT_TEXT
    p.font.name = 'Calibri'
    p.alignment = PP_ALIGN.CENTER
    p.space_before = Pt(8)


# ════════════════════════════════════════════════════════════
# SLIDE 4 — ARQUITECTURA GENERAL
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(slide)
add_title_bar(slide, "🏗️", "Arquitectura General", BLUE)

layers = [
    ("🖥️  Frontend", "React 19 + Vite + Mapbox GL", BLUE, Inches(4.5)),
    ("⬇️  REST API  ·  puerto 3333", "", DARK_GRAY, Inches(0.5)),
    ("⚙️  Backend", "NestJS + Clean Architecture + TypeScript", GREEN, Inches(4.5)),
    ("⬇️  Prisma ORM", "", DARK_GRAY, Inches(0.5)),
    ("🗄️  Base de Datos", "PostgreSQL + PostGIS", RGBColor(0x47, 0x56, 0x69), Inches(4.5)),
]

y_start = Inches(1.4)
for i, (title, subtitle, color, h) in enumerate(layers):
    y = y_start + i * Inches(1.1)
    if i % 2 == 0:  # layer block
        shape = add_shape(slide, Inches(3.5), y, Inches(6), h,
                          fill_color=CARD_BG, line_color=color)
        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_left = Pt(16)
        tf.margin_top = Pt(8) if h > Inches(1) else Pt(4)
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = color
        p.font.name = 'Calibri'
        if subtitle:
            p = tf.add_paragraph()
            p.text = subtitle
            p.font.size = Pt(13)
            p.font.color.rgb = GRAY
            p.font.name = 'Calibri'
    else:  # arrow
        add_text_box(slide, Inches(4.5), y, Inches(4), h,
                     title, font_size=11, color=DARK_GRAY,
                     alignment=PP_ALIGN.CENTER)

# Connection note
add_text_box(slide, Inches(0.8), Inches(6.5), Inches(11), Inches(0.5),
             "🔄  Comunicación mediante JSON con geometrías GeoJSON",
             font_size=13, color=DARK_GRAY, alignment=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════
# SLIDE 5 — BASE DE DATOS Y MODELO ESPACIAL
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(slide)
add_title_bar(slide, "🗃️", "Base de Datos & Modelo Espacial", GREEN)

entities = [
    ("🚆", "Líneas", "Ruta de transporte\ncon color y estado", BLUE),
    ("📍", "Estaciones", "Polígono\n(área geográfica)", GREEN),
    ("🔌", "Cables", "LineString\n(trazado de cable)", BLUE),
    ("⚡", "Postes", "Point\n(ubicación exacta)", GREEN),
    ("🔗", "Linea-Estación", "Join table\nM:N", BLUE),
]
ent_w = Inches(2.1)
ent_h = Inches(3.0)
ent_gap = Inches(0.25)
ent_total = 5 * ent_w + 4 * ent_gap
ent_start = (SLIDE_W - ent_total) // 2

for i, (icon, name, desc, color) in enumerate(entities):
    x = ent_start + i * (ent_w + ent_gap)
    shape = add_shape(slide, x, Inches(1.6), ent_w, ent_h,
                      fill_color=CARD_BG, line_color=CARD_BORDER)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(8)
    tf.margin_right = Pt(8)
    tf.margin_top = Pt(14)
    p = tf.paragraphs[0]
    p.text = icon
    p.font.size = Pt(28)
    p.alignment = PP_ALIGN.CENTER
    p = tf.add_paragraph()
    p.text = name
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = color
    p.font.name = 'Calibri'
    p.alignment = PP_ALIGN.CENTER
    p.space_before = Pt(6)
    p = tf.add_paragraph()
    p.text = desc
    p.font.size = Pt(11)
    p.font.color.rgb = GRAY
    p.font.name = 'Calibri'
    p.alignment = PP_ALIGN.CENTER
    p.space_before = Pt(4)

# Badges
badges = [
    ("🌍  PostGIS · geometry", BLUE),
    ("📦  Prisma ORM", GREEN),
    ("📤  Exportación GeoJSON", BLUE),
]
for i, (text, color) in enumerate(badges):
    x = Inches(2.0) + i * Inches(3.5)
    shape = add_shape(slide, x, Inches(5.2), Inches(2.8), Inches(0.6),
                      fill_color=RGBColor(0x14, 0x25, 0x45), line_color=color)
    tf = shape.text_frame
    tf.margin_top = Pt(4)
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(12)
    p.font.color.rgb = color
    p.font.name = 'Calibri'
    p.alignment = PP_ALIGN.CENTER

add_text_box(slide, Inches(0.8), Inches(6.2), Inches(11), Inches(0.5),
             "🔗  Una Línea tiene muchos Cables, Postes y Estaciones (a través de LineaEstación)",
             font_size=13, color=DARK_GRAY, alignment=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════
# SLIDE 6 — BACKEND: CLEAN ARCHITECTURE
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(slide)
add_title_bar(slide, "🧱", "Backend: Clean Architecture", BLUE)

layers_data = [
    ("📦 Dominio", "Entidades y reglas de negocio", BLUE),
    ("⚙️  Aplicación", "Casos de uso, DTOs y comandos", GREEN),
    ("🔧 Infraestructura", "Prisma, controladores, módulos", BLUE),
]
lay_w = Inches(3.2)
lay_h = Inches(3.0)
lay_gap = Inches(0.4)
lay_total = 3 * lay_w + 2 * lay_gap
lay_start = (SLIDE_W - lay_total) // 2

for i, (title, desc, color) in enumerate(layers_data):
    x = lay_start + i * (lay_w + lay_gap)
    shape = add_shape(slide, x, Inches(1.8), lay_w, lay_h,
                      fill_color=CARD_BG, line_color=color)
    # Left accent bar
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, x, Inches(1.8), Pt(6), lay_h
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()

    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(20)
    tf.margin_right = Pt(12)
    tf.margin_top = Pt(20)
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = color
    p.font.name = 'Calibri'
    p = tf.add_paragraph()
    p.text = desc
    p.font.size = Pt(14)
    p.font.color.rgb = LIGHT_TEXT
    p.font.name = 'Calibri'
    p.space_before = Pt(10)

# Tech badges
techs = [("🟦  TypeScript", BLUE), ("🪺  NestJS", GREEN), ("🧩  Modular por entidad", BLUE)]
for i, (text, color) in enumerate(techs):
    x = Inches(2.0) + i * Inches(3.5)
    shape = add_shape(slide, x, Inches(5.4), Inches(2.8), Inches(0.6),
                      fill_color=RGBColor(0x14, 0x25, 0x45), line_color=color)
    tf = shape.text_frame
    tf.margin_top = Pt(4)
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(12)
    p.font.color.rgb = color
    p.font.name = 'Calibri'
    p.alignment = PP_ALIGN.CENTER


# ════════════════════════════════════════════════════════════
# SLIDE 7 — API REST
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(slide)
add_title_bar(slide, "🌐", "API REST — Endpoints", GREEN)

endpoints = [
    ("GET · POST", "/api/lineas", "Líneas de transporte", BLUE),
    ("GET · POST", "/api/estaciones", "Estaciones (polígono)", GREEN),
    ("GET · POST", "/cable", "Cables (LineString)", BLUE),
    ("GET · POST", "/api/poste", "Postes (Point)", GREEN),
]

ep_w = Inches(2.6)
ep_h = Inches(2.8)
ep_gap = Inches(0.35)
ep_total = 4 * ep_w + 3 * ep_gap
ep_start = (SLIDE_W - ep_total) // 2

for i, (method, route, desc, color) in enumerate(endpoints):
    x = ep_start + i * (ep_w + ep_gap)
    shape = add_shape(slide, x, Inches(1.8), ep_w, ep_h,
                      fill_color=CARD_BG, line_color=CARD_BORDER)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(10)
    tf.margin_right = Pt(10)
    tf.margin_top = Pt(14)

    # Method badge
    p = tf.paragraphs[0]
    p.text = method
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = color
    p.font.name = 'Calibri'
    p.alignment = PP_ALIGN.CENTER

    # Route
    p = tf.add_paragraph()
    p.text = route
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = 'Consolas'
    p.alignment = PP_ALIGN.CENTER
    p.space_before = Pt(10)

    # Description
    p = tf.add_paragraph()
    p.text = desc
    p.font.size = Pt(12)
    p.font.color.rgb = GRAY
    p.font.name = 'Calibri'
    p.alignment = PP_ALIGN.CENTER
    p.space_before = Pt(6)

# Bottom info
add_text_box(slide, Inches(0.8), Inches(5.2), Inches(5), Inches(0.4),
             "📐  Respuestas con geometrías en formato GeoJSON",
             font_size=13, color=GRAY)
add_text_box(slide, Inches(6.5), Inches(5.2), Inches(5), Inches(0.4),
             "🔒  Autenticación planeada: JWT + bcrypt + cookies",
             font_size=13, color=GRAY)


# ════════════════════════════════════════════════════════════
# SLIDE 8 — FRONTEND: MAPA INTERACTIVO 3D
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(slide)
add_title_bar(slide, "🗺️", "Frontend: Mapa Interactivo 3D", BLUE)

# Tech cards
front_techs = [
    ("⚛️", "React 19", "+ Vite + TypeScript", BLUE),
    ("🗺️", "Mapbox GL", "Renderizado 3D", GREEN),
    ("🎛️", "Panel de Capas", "Toggle estaciones,\npostes, cables, 3D", BLUE),
    ("💬", "Popups Hover", "Información\nal pasar el mouse", GREEN),
]
ft_w = Inches(2.5)
ft_h = Inches(3.0)
ft_gap = Inches(0.3)
ft_total = 4 * ft_w + 3 * ft_gap
ft_start = (SLIDE_W - ft_total) // 2

for i, (icon, title, desc, color) in enumerate(front_techs):
    x = ft_start + i * (ft_w + ft_gap)
    shape = add_shape(slide, x, Inches(1.6), ft_w, ft_h,
                      fill_color=CARD_BG, line_color=CARD_BORDER)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(8)
    tf.margin_right = Pt(8)
    tf.margin_top = Pt(16)
    p = tf.paragraphs[0]
    p.text = icon
    p.font.size = Pt(30)
    p.alignment = PP_ALIGN.CENTER
    p = tf.add_paragraph()
    p.text = title
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = color
    p.font.name = 'Calibri'
    p.alignment = PP_ALIGN.CENTER
    p.space_before = Pt(6)
    p = tf.add_paragraph()
    p.text = desc
    p.font.size = Pt(12)
    p.font.color.rgb = GRAY
    p.font.name = 'Calibri'
    p.alignment = PP_ALIGN.CENTER
    p.space_before = Pt(4)

# Bottom layer types
layer_info = [
    ("■  Estaciones", " polígonos 3D", BLUE),
    ("◆  Postes", " marcadores", GREEN),
    ("━  Cables", " líneas", BLUE),
    ("🏢  Edificios", " 3D", GREEN),
]
for i, (label, desc, color) in enumerate(layer_info):
    x = Inches(1.2) + i * Inches(2.8)
    add_text_box(slide, x, Inches(5.2), Inches(2.6), Inches(0.4),
                 f"{label}{desc}", font_size=12, color=color)


# ════════════════════════════════════════════════════════════
# SLIDE 9 — DEMOSTRACIÓN
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(slide)
add_title_bar(slide, "▶️", "Demostración en Vivo", GREEN)

# Map simulation box
map_box = add_shape(slide, Inches(2.0), Inches(1.5), Inches(9), Inches(4.5),
                    fill_color=RGBColor(0x1A, 0x26, 0x42),
                    line_color=CARD_BORDER)
tf = map_box.text_frame
tf.word_wrap = True
tf.margin_top = Pt(40)
p = tf.paragraphs[0]
p.text = "🗺️"
p.font.size = Pt(64)
p.alignment = PP_ALIGN.CENTER
p = tf.add_paragraph()
p.text = "Visualización de la red en mapa 3D"
p.font.size = Pt(20)
p.font.color.rgb = GRAY
p.font.name = 'Calibri'
p.alignment = PP_ALIGN.CENTER
p.space_before = Pt(12)
p = tf.add_paragraph()
p.text = "📍 Mapbox GL"
p.font.size = Pt(12)
p.font.color.rgb = DARK_GRAY
p.font.name = 'Calibri'
p.alignment = PP_ALIGN.CENTER

add_text_box(slide, Inches(2.0), Inches(6.3), Inches(9), Inches(0.5),
             "🖥️  Demo en vivo — navegación por el mapa interactivo con estaciones, postes y cables",
             font_size=14, color=DARK_GRAY, alignment=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════
# SLIDE 10 — STACK TECNOLÓGICO
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(slide)
add_title_bar(slide, "🛠️", "Stack Tecnológico", BLUE)

# Backend section
add_text_box(slide, Inches(0.8), Inches(1.3), Inches(5), Inches(0.5),
             "🖥️  Backend", font_size=18, color=WHITE, bold=True)

back_techs = [("🟦", "NestJS"), ("📘", "TypeScript"), ("📦", "Prisma"),
              ("🗄️", "PostgreSQL"), ("🌍", "PostGIS")]
for i, (icon, name) in enumerate(back_techs):
    x = Inches(1.0) + i * Inches(2.3)
    shape = add_shape(slide, x, Inches(1.9), Inches(2.0), Inches(1.2),
                      fill_color=CARD_BG, line_color=CARD_BORDER)
    tf = shape.text_frame
    tf.margin_top = Pt(8)
    p = tf.paragraphs[0]
    p.text = f"{icon}  {name}"
    p.font.size = Pt(14)
    p.font.color.rgb = LIGHT_TEXT
    p.font.name = 'Calibri'
    p.alignment = PP_ALIGN.CENTER

# Frontend section
add_text_box(slide, Inches(0.8), Inches(3.6), Inches(5), Inches(0.5),
             "🎨  Frontend", font_size=18, color=WHITE, bold=True)

front_techs = [("⚛️", "React 19"), ("⚡", "Vite"), ("🗺️", "Mapbox GL"), ("📘", "TypeScript")]
for i, (icon, name) in enumerate(front_techs):
    x = Inches(1.0) + i * Inches(2.8)
    shape = add_shape(slide, x, Inches(4.2), Inches(2.5), Inches(1.2),
                      fill_color=CARD_BG, line_color=CARD_BORDER)
    tf = shape.text_frame
    tf.margin_top = Pt(8)
    p = tf.paragraphs[0]
    p.text = f"{icon}  {name}"
    p.font.size = Pt(14)
    p.font.color.rgb = LIGHT_TEXT
    p.font.name = 'Calibri'
    p.alignment = PP_ALIGN.CENTER

# Tools section
add_text_box(slide, Inches(0.8), Inches(5.8), Inches(5), Inches(0.5),
             "🔧  Herramientas", font_size=18, color=WHITE, bold=True)

tools = [("🐳", "Docker"), ("🔒", "JWT"), ("🛡️", "bcrypt"), ("📧", "nodemailer")]
for i, (icon, name) in enumerate(tools):
    x = Inches(1.0) + i * Inches(2.8)
    shape = add_shape(slide, x, Inches(6.3), Inches(2.5), Inches(0.8),
                      fill_color=RGBColor(0x14, 0x25, 0x45),
                      line_color=CARD_BORDER)
    tf = shape.text_frame
    tf.margin_top = Pt(4)
    p = tf.paragraphs[0]
    p.text = f"{icon}  {name}"
    p.font.size = Pt(13)
    p.font.color.rgb = LIGHT_TEXT
    p.font.name = 'Calibri'
    p.alignment = PP_ALIGN.CENTER


# ════════════════════════════════════════════════════════════
# SLIDE 11 — CONCLUSIONES Y TRABAJO FUTURO
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(slide)
add_title_bar(slide, "🏁", "Conclusiones & Trabajo Futuro", GREEN)

# Achievements
shape = add_shape(slide, Inches(1.2), Inches(1.8), Inches(5.0), Inches(4.0),
                  fill_color=CARD_BG, line_color=GREEN)
# Left accent
bar = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(1.2), Inches(1.8), Pt(6), Inches(4.0)
)
bar.fill.solid()
bar.fill.fore_color.rgb = GREEN
bar.line.fill.background()

tf = shape.text_frame
tf.word_wrap = True
tf.margin_left = Pt(24)
tf.margin_right = Pt(12)
tf.margin_top = Pt(16)
p = tf.paragraphs[0]
p.text = "✅  Logros"
p.font.size = Pt(22)
p.font.bold = True
p.font.color.rgb = GREEN
p.font.name = 'Calibri'
achievements = [
    "Visor SIG funcional con datos espaciales reales",
    "API REST con operaciones CRUD completas",
    "Arquitectura limpia y mantenible (Clean Architecture)",
    "Visualización 3D con Mapbox GL",
]
for a in achievements:
    p = tf.add_paragraph()
    p.text = f"  •  {a}"
    p.font.size = Pt(14)
    p.font.color.rgb = LIGHT_TEXT
    p.font.name = 'Calibri'
    p.space_before = Pt(8)

# Future work
shape2 = add_shape(slide, Inches(6.8), Inches(1.8), Inches(5.0), Inches(4.0),
                   fill_color=CARD_BG, line_color=BLUE)
bar2 = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(6.8), Inches(1.8), Pt(6), Inches(4.0)
)
bar2.fill.solid()
bar2.fill.fore_color.rgb = BLUE
bar2.line.fill.background()

tf2 = shape2.text_frame
tf2.word_wrap = True
tf2.margin_left = Pt(24)
tf2.margin_right = Pt(12)
tf2.margin_top = Pt(16)
p = tf2.paragraphs[0]
p.text = "🚀  Trabajo Futuro"
p.font.size = Pt(22)
p.font.bold = True
p.font.color.rgb = BLUE
p.font.name = 'Calibri'
future = [
    "Autenticación completa (JWT, roles)",
    "Panel de administración web",
    "Aplicación móvil",
    "Análisis de rutas y reportes",
]
for f in future:
    p = tf2.add_paragraph()
    p.text = f"  •  {f}"
    p.font.size = Pt(14)
    p.font.color.rgb = LIGHT_TEXT
    p.font.name = 'Calibri'
    p.space_before = Pt(8)


# ════════════════════════════════════════════════════════════
# SLIDE 12 — PREGUNTAS
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(slide)

# Gradient bar effect
bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(3.0), SLIDE_W, Inches(2.5))
bar.fill.solid()
bar.fill.fore_color.rgb = RGBColor(0x12, 0x1D, 0x3A)
bar.line.fill.background()

add_text_box(slide, Inches(2), Inches(1.2), Inches(9), Inches(1.2),
             "¡Gracias!", font_size=48, color=WHITE, bold=True,
             alignment=PP_ALIGN.CENTER)

add_divider(slide, Inches(5.5), Inches(2.4), Inches(2), GREEN)

add_text_box(slide, Inches(2), Inches(2.8), Inches(9), Inches(0.8),
             "¿Preguntas?", font_size=28, color=GRAY, bold=False,
             alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(2), Inches(4.2), Inches(9), Inches(0.5),
             "📧  contacto@universidad.edu      🐙  github.com/universidad/sig-teleferico",
             font_size=14, color=DARK_GRAY, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(2), Inches(5.5), Inches(9), Inches(0.4),
             "SIG Teleférico Universitario  ·  2026",
             font_size=11, color=RGBColor(0x33, 0x41, 0x55), alignment=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════
# GUARDAR
# ════════════════════════════════════════════════════════════
output_path = "/home/nicolas/Desktop/SIG-II/presentacion.pptx"
prs.save(output_path)
print(f"✅  Presentación guardada en: {output_path}")
